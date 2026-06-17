# ta/tools/office.py
import json
import re
from pathlib import Path

from langchain_core.tools import tool


def _add_runs_with_bold(paragraph, text: str) -> None:
    """Add text to a docx paragraph, rendering **bold** spans as bold runs."""
    for i, part in enumerate(text.split("**")):
        if not part:
            continue
        run = paragraph.add_run(part)
        if i % 2 == 1:
            run.bold = True


def _render_markdown_to_doc(doc, content: str) -> None:
    """Render a markdown-ish string into a python-docx Document.
    Supports: # .. #### headings, - / * bullets, '1.' numbered lists, **bold**
    inline, fenced ``` code blocks, and | pipe | tables |."""
    from docx.shared import Pt

    def add_code(lines: list[str]) -> None:
        if not lines:
            return
        p = doc.add_paragraph()
        run = p.add_run("\n".join(lines))
        run.font.name = "Consolas"
        run.font.size = Pt(9)

    def add_table(rows_raw: list[str]) -> None:
        rows = [[c.strip() for c in r.strip().strip("|").split("|")] for r in rows_raw]
        rows = [r for r in rows if not all(c and set(c) <= set("-: ") for c in r)]
        if not rows:
            return
        ncols = max(len(r) for r in rows)
        table = doc.add_table(rows=0, cols=ncols)
        table.style = "Table Grid"
        for r in rows:
            cells = table.add_row().cells
            for ci in range(ncols):
                cells[ci].text = r[ci] if ci < len(r) else ""

    in_code = False
    code_buf: list[str] = []
    table_buf: list[str] = []
    for raw in content.split("\n"):
        s = raw.strip()
        if s.startswith("```"):
            if in_code:
                add_code(code_buf)
                code_buf, in_code = [], False
            else:
                add_table(table_buf)
                table_buf = []
                in_code = True
            continue
        if in_code:
            code_buf.append(raw)
            continue
        if s.startswith("|") and s.endswith("|") and len(s) > 1:
            table_buf.append(s)
            continue
        if table_buf:
            add_table(table_buf)
            table_buf = []
        if not s:
            doc.add_paragraph("")
        elif s.startswith("#### "):
            doc.add_heading(s[5:], level=4)
        elif s.startswith("### "):
            doc.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s[:2] in ("- ", "* "):
            _add_runs_with_bold(doc.add_paragraph(style="List Bullet"), s[2:])
        elif re.match(r"^\d+\.\s", s):
            _add_runs_with_bold(doc.add_paragraph(style="List Number"), re.sub(r"^\d+\.\s", "", s))
        else:
            _add_runs_with_bold(doc.add_paragraph(), s)
    if table_buf:
        add_table(table_buf)
    if in_code:
        add_code(code_buf)


@tool
def read_word_file(file_path: str) -> str:
    """Read the full text content of a local Word (.docx) file, including tables."""
    try:
        from docx import Document
        path = Path(file_path)
        if not path.exists():
            return f"ERROR: File not found at '{file_path}'"
        if path.suffix.lower() != ".docx":
            return f"ERROR: Expected a .docx file, got '{path.suffix}'"
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        table_texts = []
        for i, table in enumerate(doc.tables):
            table_texts.append(f"\n[Table {i + 1}]")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_texts.append(row_text)
        content = "\n".join(paragraphs)
        if table_texts:
            content += "\n" + "\n".join(table_texts)
        return content if content.strip() else "[Document is empty]"
    except ImportError:
        return "ERROR: python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        return f"ERROR reading Word file: {e}"


@tool
def write_word_file(file_path: str, content: str, overwrite: bool = True) -> str:
    """Create or overwrite a Word (.docx) file from markdown-style text.
    Supports '#'..'####' headings, '- '/'* ' bullets, '1.' numbered lists,
    '**bold**' inline, fenced ``` code blocks, and | pipe | tables |.
    Set overwrite=False to refuse clobbering an existing file (protects instructor edits)."""
    try:
        from docx import Document
        path = Path(file_path)
        if path.exists() and not overwrite:
            return (
                f"REFUSED: '{path.resolve()}' already exists. Call again with "
                "overwrite=True to replace it, or write to a new path."
            )
        existed = path.exists()
        path.parent.mkdir(parents=True, exist_ok=True)
        doc = Document()
        _render_markdown_to_doc(doc, content)
        doc.save(str(path))
        note = " (replaced existing file)" if existed else ""
        return f"SUCCESS: Word document saved to '{path.resolve()}'{note}"
    except ImportError:
        return "ERROR: python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        return f"ERROR writing Word file: {e}"


@tool
def append_to_word_file(file_path: str, content: str) -> str:
    """Append markdown-style content to an existing Word (.docx) file. Creates it if absent.
    Same rich markdown syntax as write_word_file (headings, bold, lists, code, tables)."""
    try:
        from docx import Document
        path = Path(file_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        doc = Document(str(path)) if path.exists() else Document()
        _render_markdown_to_doc(doc, content)
        doc.save(str(path))
        return f"SUCCESS: Content appended to '{path.resolve()}'"
    except ImportError:
        return "ERROR: python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        return f"ERROR appending to Word file: {e}"


@tool
def read_excel_file(file_path: str, sheet_name: str = "") -> str:
    """Read data from a local Excel (.xlsx) file. sheet_name: specific sheet or empty for all."""
    try:
        import pandas as pd
        path = Path(file_path)
        if not path.exists():
            return f"ERROR: File not found at '{file_path}'"
        if path.suffix.lower() not in (".xlsx", ".xls", ".xlsm"):
            return f"ERROR: Expected an Excel file, got '{path.suffix}'"
        xl = pd.ExcelFile(str(path))
        sheets_to_read = (
            xl.sheet_names if not sheet_name or sheet_name.upper() == "ALL"
            else [sheet_name] if sheet_name in xl.sheet_names
            else None
        )
        if sheets_to_read is None:
            return f"ERROR: Sheet '{sheet_name}' not found. Available: {xl.sheet_names}"
        parts = []
        for sname in sheets_to_read:
            df = pd.read_excel(str(path), sheet_name=sname)
            parts.append(f"[Sheet: {sname}] ({len(df)} rows x {len(df.columns)} cols)")
            parts.append(df.to_string(index=False))
            parts.append("")
        return "\n".join(parts)
    except ImportError:
        return "ERROR: pandas or openpyxl not installed. Run: pip install pandas openpyxl"
    except Exception as e:
        return f"ERROR reading Excel file: {e}"


@tool
def get_excel_sheet_names(file_path: str) -> str:
    """List all sheet names in a local Excel (.xlsx) file."""
    try:
        import openpyxl
        path = Path(file_path)
        if not path.exists():
            return f"ERROR: File not found at '{file_path}'"
        wb = openpyxl.load_workbook(str(path), read_only=True)
        sheets = wb.sheetnames
        wb.close()
        return f"Sheets found: {', '.join(sheets)}"
    except ImportError:
        return "ERROR: openpyxl not installed. Run: pip install openpyxl"
    except Exception as e:
        return f"ERROR reading Excel file: {e}"


@tool
def write_excel_file(
    file_path: str, data_json: str, sheet_name: str = "Sheet1", overwrite: bool = True
) -> str:
    """Create or overwrite a local Excel (.xlsx) file from a JSON array of row dicts.
    Example data_json: '[{"Name": "Alice", "Score": 95}]'
    Set overwrite=False to refuse clobbering an existing file."""
    try:
        import pandas as pd
        path = Path(file_path)
        if path.exists() and not overwrite:
            return (
                f"REFUSED: '{path.resolve()}' already exists. Call again with "
                "overwrite=True to replace it, or write to a new path."
            )
        path.parent.mkdir(parents=True, exist_ok=True)
        try:
            data = json.loads(data_json)
        except json.JSONDecodeError as je:
            return f"ERROR: Invalid JSON: {je}"
        if not isinstance(data, list):
            return "ERROR: data_json must be a JSON array."
        df = pd.DataFrame(data)
        df.to_excel(str(path), index=False, sheet_name=sheet_name)
        return (
            f"SUCCESS: Excel saved to '{path.resolve()}' — "
            f"{len(df)} rows, {len(df.columns)} cols, sheet '{sheet_name}'"
        )
    except ImportError:
        return "ERROR: pandas or openpyxl not installed. Run: pip install pandas openpyxl"
    except Exception as e:
        return f"ERROR writing Excel file: {e}"


@tool
def append_excel_rows(file_path: str, data_json: str, sheet_name: str = "Sheet1") -> str:
    """Append rows to a local Excel (.xlsx) sheet. Creates file/sheet if absent.
    data_json: JSON array of row dicts."""
    try:
        import pandas as pd
        path = Path(file_path)
        try:
            new_data = json.loads(data_json)
        except json.JSONDecodeError as je:
            return f"ERROR: Invalid JSON: {je}"
        if not isinstance(new_data, list):
            return "ERROR: data_json must be a JSON array."
        new_df = pd.DataFrame(new_data)
        if path.exists():
            try:
                xl = pd.ExcelFile(str(path))
                sheet_exists = sheet_name in xl.sheet_names
            except Exception:
                sheet_exists = False

            if sheet_exists:
                with pd.ExcelWriter(str(path), engine="openpyxl", mode="a",
                                    if_sheet_exists="overlay") as writer:
                    existing_df = pd.read_excel(str(path), sheet_name=sheet_name)
                    pd.concat([existing_df, new_df], ignore_index=True).to_excel(
                        writer, index=False, sheet_name=sheet_name
                    )
            else:
                with pd.ExcelWriter(str(path), engine="openpyxl", mode="a") as writer:
                    new_df.to_excel(writer, index=False, sheet_name=sheet_name)
        else:
            path.parent.mkdir(parents=True, exist_ok=True)
            new_df.to_excel(str(path), index=False, sheet_name=sheet_name)
        return (
            f"SUCCESS: {len(new_data)} row(s) appended to "
            f"sheet '{sheet_name}' in '{path.resolve()}'"
        )
    except ImportError:
        return "ERROR: pandas or openpyxl not installed. Run: pip install pandas openpyxl"
    except Exception as e:
        return f"ERROR appending to Excel file: {e}"


@tool
def read_pptx_file(file_path: str) -> str:
    """Read and extract all text from a local PowerPoint (.pptx) file, organized by slide."""
    try:
        from pptx import Presentation
        path = Path(file_path)
        if not path.exists():
            return f"ERROR: File not found at '{file_path}'"
        if path.suffix.lower() != ".pptx":
            return f"ERROR: Expected a .pptx file, got '{path.suffix}'"
        prs = Presentation(str(path))
        output = []
        for i, slide in enumerate(prs.slides, 1):
            output.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            output.append(text)
            output.append("")
        return "\n".join(output) if output else "[Presentation has no slides]"
    except ImportError:
        return "ERROR: python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"ERROR reading PowerPoint file: {e}"


@tool
def write_pptx_file(file_path: str, slides_json: str, overwrite: bool = True) -> str:
    """Create a local PowerPoint (.pptx) from a JSON array of slide objects.
    Each slide: {"title": str, "content": [str, ...], "notes": str (optional)}.
    Example: '[{"title": "Intro", "content": ["Point 1", "Point 2"]}]'
    Set overwrite=False to refuse clobbering an existing file."""
    try:
        from pptx import Presentation
        from pptx.util import Pt  # noqa: F401
        path = Path(file_path)
        if path.exists() and not overwrite:
            return (
                f"REFUSED: '{path.resolve()}' already exists. Call again with "
                "overwrite=True to replace it, or write to a new path."
            )
        path.parent.mkdir(parents=True, exist_ok=True)
        try:
            slides_data = json.loads(slides_json)
        except json.JSONDecodeError as je:
            return f"ERROR: Invalid JSON: {je}"
        if not isinstance(slides_data, list):
            return "ERROR: slides_json must be a JSON array."
        prs = Presentation()
        for i, slide_data in enumerate(slides_data):
            if not isinstance(slide_data, dict):
                continue
            layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("title", f"Slide {i + 1}")
            content_items = slide_data.get("content", [])
            if content_items and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for j, item in enumerate(content_items):
                    if j == 0:
                        tf.text = item
                    else:
                        p = tf.add_paragraph()
                        p.text = item
            if slide_data.get("notes"):
                slide.notes_slide.notes_text_frame.text = slide_data["notes"]
        prs.save(str(path))
        return f"SUCCESS: PowerPoint saved to '{path.resolve()}' with {len(slides_data)} slide(s)"
    except ImportError:
        return "ERROR: python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"ERROR writing PowerPoint file: {e}"


@tool
def list_office_files(directory: str, extension: str = "") -> str:
    """List Office files (.docx, .xlsx, .pptx) in a directory.
    extension: optional filter — 'docx', 'xlsx', or 'pptx'. Empty = all Office files."""
    try:
        path = Path(directory)
        if not path.exists():
            return f"ERROR: Directory not found at '{directory}'"
        if not path.is_dir():
            return f"ERROR: '{directory}' is not a directory."
        office_extensions = {".docx", ".xlsx", ".xls", ".xlsm", ".pptx"}
        target_exts = {f".{extension.lower().lstrip('.')}"} if extension else office_extensions
        found = sorted(
            str(f.resolve()) for f in path.rglob("*") if f.suffix.lower() in target_exts
        )
        if not found:
            return f"No Office files found in '{directory}'"
        return "Found files:\n" + "\n".join(f"  - {f}" for f in found)
    except Exception as e:
        return f"ERROR listing files: {e}"


def _pdf_safe(s: str) -> str:
    """Map common Unicode punctuation to ASCII, then drop anything outside latin-1
    so fpdf2 core fonts (which are latin-1) never crash. Spanish accents survive."""
    for k, v in {
        "‘": "'", "’": "'", "“": '"', "”": '"',
        "–": "-", "—": "-", "…": "...", "•": "-",
    }.items():
        s = s.replace(k, v)
    return s.encode("latin-1", "replace").decode("latin-1")


def _render_markdown_to_pdf(content: str, out_path: Path) -> None:
    """Render markdown-ish text to a PDF via fpdf2 (pure-python, no Word).
    Headings, bullets, numbered lists, and fenced code blocks; **bold** markers stripped."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    def mc(h: float, txt: str) -> None:
        # new_x=LMARGIN resets the cursor to the left margin so the next line has
        # full page width (default leaves x at the right margin -> zero-width error).
        pdf.multi_cell(0, h, _pdf_safe(txt), new_x="LMARGIN", new_y="NEXT")

    in_code = False
    for raw in content.split("\n"):
        s = raw.strip()
        if s.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            pdf.set_font("Courier", size=9)
            mc(5, raw)
            continue
        if not s:
            pdf.ln(3)
            continue
        if s.startswith("#### "):
            pdf.set_font("Helvetica", "B", 11)
            mc(7, s[5:])
        elif s.startswith("### "):
            pdf.set_font("Helvetica", "B", 12)
            mc(8, s[4:])
        elif s.startswith("## "):
            pdf.set_font("Helvetica", "B", 14)
            mc(9, s[3:])
        elif s.startswith("# "):
            pdf.set_font("Helvetica", "B", 16)
            mc(10, s[2:])
        elif s[:2] in ("- ", "* "):
            pdf.set_font("Helvetica", size=11)
            mc(6, "  - " + s[2:].replace("**", ""))
        elif re.match(r"^\d+\.\s", s):
            pdf.set_font("Helvetica", size=11)
            mc(6, "  " + s.replace("**", ""))
        else:
            pdf.set_font("Helvetica", size=11)
            mc(6, s.replace("**", ""))
    pdf.output(str(out_path))


@tool
def export_to_pdf(source_path: str, output_path: str = "") -> str:
    """Export a document to PDF for student distribution. Accepts a .md/.markdown/.txt
    file or a .docx file. Pure-python rendering (headings, bullets, numbered lists, code
    blocks) — no Microsoft Word needed. output_path defaults to the source name with .pdf."""
    try:
        src = Path(source_path)
        if not src.exists():
            return f"ERROR: File not found at '{source_path}'"
        suffix = src.suffix.lower()
        if suffix in (".md", ".markdown", ".txt"):
            text = src.read_text(encoding="utf-8")
        elif suffix == ".docx":
            text = read_word_file.func(str(src))
            if text.startswith("ERROR"):
                return text
        else:
            return f"ERROR: Unsupported source '{src.suffix}'. Use .md, .txt, or .docx."

        out = Path(output_path) if output_path else src.with_suffix(".pdf")
        out.parent.mkdir(parents=True, exist_ok=True)
        try:
            _render_markdown_to_pdf(text, out)
        except ImportError:
            return "ERROR: fpdf2 not installed. Run: uv add fpdf2"
        return f"SUCCESS: PDF saved to '{out.resolve()}'"
    except Exception as e:
        return f"ERROR exporting to PDF: {e}"
