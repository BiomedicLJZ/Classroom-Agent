# ta/tools/drive.py
import io
from functools import cache

from googleapiclient.discovery import build
from langchain_core.tools import tool

from ta.google_auth import get_credentials
from ta.session import get_active_account

_GOOGLE_DOC_MIMES = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}
_TEXT_MIMES = {
    "text/plain", "text/x-python", "text/markdown",
    "text/x-markdown", "application/json", "text/html",
}
_OFFICE_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}


def _parse_office_bytes(fmt: str, data: bytes) -> str:
    buf = io.BytesIO(data)
    if fmt == "docx":
        from docx import Document
        doc = Document(buf)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        table_texts = []
        for i, table in enumerate(doc.tables):
            table_texts.append(f"\n[Table {i + 1}]")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_texts.append(row_text)
        content = "\n".join(paragraphs)
        if table_texts:
            content += "\n" + "\n".join(table_texts)
        return content or "[Document is empty]"
    if fmt == "xlsx":
        import pandas as pd
        xl = pd.ExcelFile(buf)
        parts = []
        for sname in xl.sheet_names:
            buf.seek(0)
            df = pd.read_excel(buf, sheet_name=sname)
            parts.append(f"[Sheet: {sname}] ({len(df)} rows x {len(df.columns)} cols)")
            parts.append(df.to_string(index=False))
        return "\n".join(parts) or "[Spreadsheet is empty]"
    if fmt == "pptx":
        from pptx import Presentation
        prs = Presentation(buf)
        output = []
        for i, slide in enumerate(prs.slides, 1):
            output.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            output.append(text)
            output.append("")
        return "\n".join(output) or "[Presentation has no text]"
    return "[Unsupported office format]"


@cache
def _drive_service(alias: str):
    creds = get_credentials(alias)
    return build("drive", "v3", credentials=creds)


@tool
def get_drive_file_text(file_id: str) -> str:
    """Download and return text content of a Google Drive file.
    Supports: Google Docs/Sheets/Slides (exported as text), Python/text/markdown files,
    PDFs (first 10 pages via pypdf)."""
    svc = _drive_service(get_active_account())
    meta = svc.files().get(fileId=file_id, fields="mimeType,name").execute()
    mime = meta.get("mimeType", "")

    if mime in _GOOGLE_DOC_MIMES:
        data = svc.files().export_media(fileId=file_id, mimeType=_GOOGLE_DOC_MIMES[mime]).execute()
        return data.decode("utf-8") if isinstance(data, bytes) else str(data)

    if mime in _TEXT_MIMES or mime.startswith("text/"):
        data = svc.files().get_media(fileId=file_id).execute()
        return data.decode("utf-8") if isinstance(data, bytes) else str(data)

    if mime == "application/pdf":
        from pypdf import PdfReader
        data = svc.files().get_media(fileId=file_id).execute()
        reader = PdfReader(io.BytesIO(data))
        return "\n".join(p.extract_text() or "" for p in reader.pages[:10])

    if mime in _OFFICE_MIMES:
        data = svc.files().get_media(fileId=file_id).execute()
        return _parse_office_bytes(_OFFICE_MIMES[mime], data)

    return f"[Cannot read file type: {mime}. File name: {meta.get('name', 'unknown')}]"


@tool
def upload_file_to_drive(local_path: str, parent_folder_id: str, filename: str) -> str:
    """Upload a local file to a Google Drive folder. Returns the new file ID."""
    from googleapiclient.http import MediaFileUpload
    svc = _drive_service(get_active_account())
    result = svc.files().create(
        body={"name": filename, "parents": [parent_folder_id]},
        media_body=MediaFileUpload(local_path, resumable=True),
        fields="id",
    ).execute()
    return f"Uploaded '{filename}' to Drive (id: {result['id']})."
